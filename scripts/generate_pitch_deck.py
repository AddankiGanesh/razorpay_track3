"""Generate RevRecover pitch deck PPTX from the 5-minute scenario script."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "RevRecover_Pitch_Deck.pptx"

# RevRecover dashboard palette
BG = RGBColor(8, 12, 20)
CARD = RGBColor(20, 29, 46)
ACCENT = RGBColor(34, 211, 238)
TEXT = RGBColor(241, 245, 249)
MUTED = RGBColor(139, 156, 179)
OK = RGBColor(52, 211, 153)
WARN = RGBColor(251, 191, 36)


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: int = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"


def add_bullets(slide, left, top, width, height, items: list[str], size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)


def title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(1.2), "RevRecover", 44, ACCENT, True)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.5),
        Inches(11),
        Inches(0.8),
        "AI Revenue Recovery Operator on Razorpay",
        24,
        TEXT,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.4),
        Inches(11),
        Inches(0.6),
        "Razorpay AI Buildathon · Track 03",
        16,
        MUTED,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.8),
        Inches(11),
        Inches(0.5),
        "Detect → Diagnose → Decide → Execute → Prove ₹ recovered",
        14,
        OK,
    )


def content_slide(prs, title: str, bullets: list[str], subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.7), title, 28, ACCENT, True)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(1.15), Inches(12), Inches(0.5), subtitle, 14, MUTED)
    add_bullets(slide, Inches(0.85), Inches(1.75), Inches(11.5), Inches(5), bullets, 17)


def scenario_slide(prs, title: str, scenario: str, razorpay: str, revrecover: str, punch: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.6), title, 26, ACCENT, True)
    add_textbox(slide, Inches(0.7), Inches(1.2), Inches(12), Inches(0.5), scenario, 18, TEXT, True)
    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(0.4), "Razorpay tells us", 13, MUTED)
    add_textbox(slide, Inches(0.7), Inches(2.35), Inches(5.5), Inches(1.2), razorpay, 15, TEXT)
    add_textbox(slide, Inches(6.5), Inches(2.0), Inches(5.8), Inches(0.4), "RevRecover does", 13, OK)
    add_textbox(slide, Inches(6.5), Inches(2.35), Inches(5.8), Inches(1.5), revrecover, 15, TEXT)
    add_textbox(slide, Inches(0.7), Inches(4.5), Inches(11.5), Inches(1.2), punch, 16, WARN)


def two_col_slide(prs, title: str, left_title: str, left_items: list[str], right_title: str, right_items: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.6), title, 26, ACCENT, True)
    add_textbox(slide, Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.4), left_title, 16, MUTED, True)
    add_bullets(slide, Inches(0.7), Inches(1.75), Inches(5.5), Inches(4.5), left_items, 15)
    add_textbox(slide, Inches(6.5), Inches(1.3), Inches(5.8), Inches(0.4), right_title, 16, OK, True)
    add_bullets(slide, Inches(6.5), Inches(1.75), Inches(5.8), Inches(4.5), right_items, 15)


def closing_slide(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8), "Razorpay shows money at risk.", 28, TEXT, True)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.0),
        Inches(11),
        Inches(0.8),
        "RevRecover shows money you can still win back — and money you already won.",
        22,
        ACCENT,
        True,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.5),
        Inches(11),
        Inches(1.2),
        "We don't chase every failed payment like an ex who won't read the hint.",
        20,
        TEXT,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(4.5),
        Inches(11),
        Inches(1.0),
        "We chase the ones who want to pay — and stop when they swiped left on checkout.",
        18,
        MUTED,
    )
    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5), "Thank you · GitHub repo · Track 03", 14, OK)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)

    content_slide(
        prs,
        "11 PM. Flash sale. Money at risk.",
        [
            "₹3 lakh expected tonight — dashboard says failed.",
            "Three customers. Same word: FAILED.",
            "But three completely different stories.",
        ],
        "Hook — scenario-based (0:00–0:45)",
    )

    content_slide(
        prs,
        "Three customers — same night",
        [
            "Rahul — wrong OTP. He wants to pay. High intent.",
            "Priya — pressed back on her phone. She changed her mind.",
            "HDFC netbanking — bank is down. Nobody can pay right now.",
            "One generic SMS to all three = spam + wasted ₹ + lost Rahul.",
        ],
    )

    two_col_slide(
        prs,
        "Razorpay today vs the merchant gap",
        "What Razorpay gives ✅",
        [
            "Webhooks: payment.failed, subscription.halted, etc.",
            "Payment Links + customer notify",
            "Subscription retry schedules (T+1, T+2, T+3)",
            "Sandbox, APIs, test mode",
        ],
        "What merchants still do manually ❌",
        [
            "Read error_reason — OTP vs cancelled vs bank down",
            "Pick channel & timing — SMS vs email vs wait",
            "Stopping rules — when NOT to nudge",
            "Prove ₹ recovered after your intervention",
        ],
    )

    content_slide(
        prs,
        "Two leaks judges care about",
        [
            "B2B: ₹25,000 invoice link expired — regenerate link, but who nudges? SMS at 2 AM or account manager?",
            "Subscriptions: halted after retries — invoices may still generate but auto-charge may never run again.",
            "Track 03 bar: measured money recovered + compliant escalation + stopping rules + audit trail.",
        ],
        "Silent revenue bleed (1:00–1:45)",
    )

    content_slide(
        prs,
        "What we built — RevRecover",
        [
            "AI revenue recovery operator on Razorpay webhooks.",
            "110+ error_reason codes → specific playbook per failure.",
            "Chase vs STOP vs DELAY vs WATCH — not one generic reminder bot.",
            "SMS · Email · Voice (Hinglish) · Payment links · Human escalation for B2B.",
            "Every recovered rupee attributed to the intervention that caused it.",
        ],
    )

    content_slide(
        prs,
        "Pipeline",
        [
            "DETECT — webhooks + lab scenarios",
            "DIAGNOSE — Razorpay error taxonomy → playbook",
            "DECIDE — stopping rules, outage delay, promise-to-pay",
            "EXECUTE — payment link, SMS, email, voice",
            "AUDIT — order.paid / payment_link.paid → ₹ recovered",
        ],
    )

    scenario_slide(
        prs,
        "Live demo · Scenario 1",
        "Wrong OTP — ₹499",
        "error_reason: incorrect_otp\nCustomer-side authentication failure",
        "Playbook: retry_with_new_otp\nUrgent SMS + payment link\nHigh priority chase",
        "Pitch line: Razorpay is the thermometer. RevRecover is the doctor.",
    )

    scenario_slide(
        prs,
        "Live demo · Scenario 2",
        "Customer cancelled — ₹799",
        "error_reason: payment_cancelled\nUser left checkout (back button / closed tab)",
        "Playbook: soft_nudge_once\nOne gentle email — not 3 aggressive SMS",
        "Same failure family, opposite actions — compliance.",
    )

    scenario_slide(
        prs,
        "Live demo · Scenario 3",
        "Bank downtime — ₹999",
        "error_reason: bank_technical_error\nBank/gateway outage window",
        "Status: DELAYED — no customer spam during outage\nRetry when bank recovers",
        "Naive systems retry immediately and burn trust.",
    )

    content_slide(
        prs,
        "More scenarios in the demo",
        [
            "Subscription halted → voice + Hinglish IVR + promise-to-pay",
            "Mandate debit declined → 3-step sequencer (SMS → email → re-register) then STOP",
            "Late auth pending → WATCHING — poll only, no customer nudge",
            "B2B link expired ₹25k → regenerate link + human escalation queue",
            "14 curated lab scenarios aligned to Razorpay error catalog",
        ],
        "Differentiators (3:45–4:15)",
    )

    content_slide(
        prs,
        "The proof — measured ₹ recovered",
        [
            "Fire all scenarios → ₹33,690 at risk in demo batch",
            "Simulate / pay links → ₹2,298 recovered with attribution",
            "Recovery by category — which playbook won the money back",
            "Counterfactual: naive remind-all vs smart chase/stop/delay",
            "Whole rupees, audit trail — CFO-ready",
        ],
        "Judging bar (4:15–4:45)",
    )

    two_col_slide(
        prs,
        "Counterfactual simulator",
        "Remind everyone (naive)",
        [
            "SMS every failed case",
            "Assume ~35% recover",
            "Annoys low-intent customers",
            "Retries during bank outages",
        ],
        "RevRecover (smart)",
        [
            "Chase high-intent cases only",
            "STOP spam & compliance caps",
            "DELAY on bank/gateway outage",
            "Voice only when worth it (B2B / halted sub)",
        ],
    )

    closing_slide(prs)

    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    main()
